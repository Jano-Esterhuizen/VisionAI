import os
import base64
import io
from openai import OpenAI 
import openai
from flask import Flask, request, send_file, render_template, jsonify
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageOps
import time

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'output'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

client = OpenAI(api_key= "")

# Function to encode an image as a base64 string
def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

# Function to extract images from a PowerPoint file
def extract_images_from_pptx(pptx_path, output_dir):
    prs = Presentation(pptx_path)
    image_paths = []

    for slide_number, slide in enumerate(prs.slides):
        for shape_number, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_path = os.path.join(output_dir, f'slide_{slide_number+1}_image_{shape_number+1}.{image.ext}')
                os.makedirs(os.path.dirname(image_path), exist_ok=True)  # Ensure the directory exists
                with open(image_path, 'wb') as f:
                    f.write(image.blob)
                image_paths.append((slide_number, shape_number, image_path))
    return image_paths

# Function to describe an image using the OpenAI API
def describe_image(image_path):
    base64_image = encode_image(image_path)
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that describes images for visually impaired users."},
            {"role": "user", "content": [
                {"type": "text", "text": "Describe this image."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
            ]}
        ],
        temperature=0.0
    )

    description = response.choices[0].message.content
    return description

# Function to insert alt text into a PowerPoint file
def insert_alt_text(pptx_path, descriptions):
    prs = Presentation(pptx_path)
    
    for slide_number, shape_number, description in descriptions:
        slide = prs.slides[slide_number]
        shape = slide.shapes[shape_number]
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"Setting alt text for slide {slide_number+1}, shape {shape_number+1}: {description}")
            shape._element.xpath('./p:nvPicPr/p:cNvPr')[0].set('descr', description)
    
    modified_pptx_path = os.path.join(OUTPUT_FOLDER, 'modified_' + os.path.basename(pptx_path))
    prs.save(modified_pptx_path)
    print(f"File saved to: {modified_pptx_path}")
    return modified_pptx_path

@app.route('/')
def index():
    return render_template('upload.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)
    
    # Extract images
    images = extract_images_from_pptx(filepath, OUTPUT_FOLDER)
    
    # Generate descriptions
    descriptions = []
    total_images = len(images)
    for i, (slide_number, shape_number, image_path) in enumerate(images):
        description = describe_image(image_path)
        descriptions.append((slide_number, shape_number, description))
        # Simulate processing time for demonstration purposes
        time.sleep(0.5)
        progress = int(((i + 1) / total_images) * 100)
        print(f"Progress: {progress}%")

    # Insert alt text and save the file
    modified_filepath = insert_alt_text(filepath, descriptions)
    
    # Get the filename for the modified file
    modified_filename = os.path.basename(modified_filepath)
    return jsonify({"filename": modified_filename})

@app.route('/download/<filename>')
def download_file(filename):
    file_path = os.path.join(OUTPUT_FOLDER, filename)
    print(f"Attempting to download: {file_path}")
    return send_file(file_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
