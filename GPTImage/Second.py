import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from transformers import BlipProcessor, BlipForConditionalGeneration
from PIL import Image

# Function to extract images from a PowerPoint file
def extract_images_from_pptx(pptx_path, output_dir):
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"The file at {pptx_path} does not exist.")

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    prs = Presentation(pptx_path)
    image_paths = []

    for slide_number, slide in enumerate(prs.slides):
        for shape_number, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_path = os.path.join(output_dir, f'slide_{slide_number+1}_image_{shape_number+1}.{image.ext}')
                with open(image_path, 'wb') as f:
                    f.write(image.blob)
                image_paths.append((slide_number, shape_number, image_path))
    return image_paths

# Function to describe an image using an AI model
def describe_image(image_path):
    processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
    model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
    
    raw_image = Image.open(image_path).convert('RGB')
    inputs = processor(raw_image, return_tensors="pt")
    out = model.generate(**inputs)
    caption = processor.decode(out[0], skip_special_tokens=True)
    return caption

# Function to insert alt text into a PowerPoint file
def insert_alt_text(pptx_path, descriptions):
    prs = Presentation(pptx_path)
    
    for slide_number, shape_number, description in descriptions:
        slide = prs.slides[slide_number]
        shape = slide.shapes[shape_number]
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"Setting alt text for slide {slide_number+1}, shape {shape_number+1}: {description}")
            shape._element.xpath('./p:nvPicPr/p:cNvPr')[0].set('descr', description)
    
    modified_pptx_path = os.path.join(os.path.dirname(pptx_path), 'modified_' + os.path.basename(pptx_path))
    prs.save(modified_pptx_path)
    return modified_pptx_path

# Example usage
pptx_path = 'C:/Users/JanoEsterhuyse/Desktop/Jano Stuff/UNI/ImageDescriber/InputFiles/Test.pptx'
output_dir = 'C:/Users/JanoEsterhuyse/Desktop/Jano Stuff/UNI/ImageDescriber/OutputFiles'

# Extract images
images = extract_images_from_pptx(pptx_path, output_dir)

# Describe all images and prepare the descriptions list
descriptions = []
for slide_number, shape_number, image_path in images:
    description = describe_image(image_path)
    descriptions.append((slide_number, shape_number, description))

# Insert the alt text into the PowerPoint file for all images
modified_pptx_path = insert_alt_text(pptx_path, descriptions)
print(f"Modified PowerPoint file saved at: {modified_pptx_path}")
