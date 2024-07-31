# import os
# from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# from tkinter import Tk
# from tkinter.filedialog import askopenfilename

# def extract_images_from_pptx(pptx_path, output_dir):
#     if not os.path.exists(pptx_path):
#         raise FileNotFoundError(f"The file at {pptx_path} does not exist.")

#     if not os.path.exists(output_dir):
#         os.makedirs(output_dir)
    
#     prs = Presentation(pptx_path)
#     image_paths = []
    
#     for slide_number, slide in enumerate(prs.slides):
#         for shape_number, shape in enumerate(slide.shapes):
#             if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#                 image = shape.image
#                 image_path = os.path.join(output_dir, f'slide_{slide_number+1}_image_{shape_number+1}.{image.ext}')
#                 with open(image_path, 'wb') as f:
#                     f.write(image.blob)
#                 image_paths.append((slide_number, shape_number, image_path))
#     return image_paths

# # Example usage with interactive file selection
# Tk().withdraw()  # we don't want a full GUI, so keep the root window from appearing
# pptx_path = askopenfilename(title="Select PowerPoint file", filetypes=[("PowerPoint files", "*.pptx")])
# output_dir = 'C:/Users/JanoEsterhuyse/Desktop/Jano Stuff/UNI/ImageDescriber/OutputFiles'
# images = extract_images_from_pptx(pptx_path, output_dir)


import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images_from_pptx(pptx_path, output_dir):
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"The file at {pptx_path} does not exist.")

    # Create output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Open the PowerPoint file
    prs = Presentation(pptx_path)
    image_paths = [] # List to hold image paths
    
    # Iterate over slides
    for slide_number, slide in enumerate(prs.slides):
        # Iterate over shapes within each slide
        for shape_number, shape in enumerate(slide.shapes):
            # Check if the shape is a picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                # Define the path to save the image
                image_path = os.path.join(output_dir, f'slide_{slide_number+1}_image_{shape_number+1}.{image.ext}')
                # Write the image blob to the defined path
                with open(image_path, 'wb') as f:
                    f.write(image.blob)
                image_paths.append((slide_number, shape_number, image_path))
                
    # Return list of image paths along with slide and shape indices
    return image_paths

# Example usage
pptx_path = 'C:/Users/JanoEsterhuyse/Desktop/Jano Stuff/UNI/ImageDescriber/InputFiles/Test.pptx'
output_dir = 'C:/Users/JanoEsterhuyse/Desktop/Jano Stuff/UNI/ImageDescriber/OutputFiles'
images = extract_images_from_pptx(pptx_path, output_dir)

from transformers import BlipProcessor, BlipForConditionalGeneration
from PIL import Image

def describe_image(image_path):
    # Load the processor and model
    processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
    model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
    
    # Open the image
    raw_image = Image.open(image_path).convert('RGB')
    # Process the image for the model
    inputs = processor(raw_image, return_tensors="pt")
    # Generate the caption
    out = model.generate(**inputs)
    # Decode the caption
    caption = processor.decode(out[0], skip_special_tokens=False)
    return caption

# Example usage
image_path = 'C:/Users/JanoEsterhuyse/Desktop/Jano Stuff/UNI/ImageDescriber/OutputFiles/slide_5_image_3.jpg'
descriptions = describe_image(image_path)

def insert_alt_text(pptx_path, descriptions):
    # Open the PowerPoint file
    prs = Presentation(pptx_path)
    
    # Iterate over the descriptions
    # for slide_number, shape_number, description in descriptions:
    #     slide = prs.slides[slide_number]
    #     shape = slide.shapes[shape_number]
    #     # Check if the shape is a picture
    #     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Set the alt text
    image_path.alt_text = descriptions
    
    # Save the modified PowerPoint file
    prs.save('modified_presentation.pptx')

# Example usage
# descriptions = [(0, 0, 'A scenic mountain view.'), (1, 1, 'A beautiful sunset over the ocean.')]
insert_alt_text(pptx_path, descriptions)










# import os
# from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# from transformers import BlipProcessor, BlipForConditionalGeneration
# from PIL import Image

# # Function to extract images from a PowerPoint file
# def extract_images_from_pptx(pptx_path, output_dir):
#     if not os.path.exists(pptx_path):
#         raise FileNotFoundError(f"The file at {pptx_path} does not exist.")

#     if not os.path.exists(output_dir):
#         os.makedirs(output_dir)
    
#     prs = Presentation(pptx_path)
#     image_paths = []

#     for slide_number, slide in enumerate(prs.slides):
#         for shape_number, shape in enumerate(slide.shapes):
#             if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#                 image = shape.image
#                 image_path = os.path.join(output_dir, f'slide_{slide_number+1}_image_{shape_number+1}.{image.ext}')
#                 with open(image_path, 'wb') as f:
#                     f.write(image.blob)
#                 image_paths.append((slide_number, shape_number, image_path))
#     return image_paths

# # Function to describe an image using an AI model
# def describe_image(image_path):
#     processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
#     model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
    
#     raw_image = Image.open(image_path).convert('RGB')
#     inputs = processor(raw_image, return_tensors="pt")
#     out = model.generate(**inputs)
#     caption = processor.decode(out[0], skip_special_tokens=True)
#     return caption

# # Function to insert alt text into a PowerPoint file
# def insert_alt_text(pptx_path, descriptions):
#     prs = Presentation(pptx_path)
    
#     for slide_number, shape_number, description in descriptions:
#         slide = prs.slides[slide_number]
#         shape = slide.shapes[shape_number]
#         if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#             shape.alt_text = description
    
#     modified_pptx_path = os.path.join(os.path.dirname(pptx_path), 'modified_' + os.path.basename(pptx_path))
#     prs.save(modified_pptx_path)
#     return modified_pptx_path

# # Main function to process the PowerPoint file
# def process_pptx(pptx_path, output_dir):
#     images = extract_images_from_pptx(pptx_path, output_dir)
#     descriptions = []
#     for slide_number, shape_number, image_path in images:
#         description = describe_image(image_path)
#         descriptions.append((slide_number, shape_number, description))
#     modified_pptx_path = insert_alt_text(pptx_path, descriptions)
#     return modified_pptx_path

# # Example usage
# pptx_path = 'C:/Users/JanoEsterhuyse/Desktop/Jano Stuff/UNI/ImageDescriber/InputFiles/Test.pptx'
# output_dir = 'C:/Users/JanoEsterhuyse/Desktop/Jano Stuff/UNI/ImageDescriber/OutputFiles'

# # Process the PowerPoint file and get the path to the modified file
# modified_pptx_path = process_pptx(pptx_path, output_dir)
# print(f"Modified PowerPoint file saved at: {modified_pptx_path}")



