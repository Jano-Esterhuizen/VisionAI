import logging
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from flask import current_app

def extract_images_from_pptx(pptx_path, output_dir):
    prs = Presentation(pptx_path)
    image_paths = []

    for slide_number, slide in enumerate(prs.slides):
        for shape_number, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_path = os.path.join(output_dir, f'slide_{slide_number+1}_image_{shape_number+1}.{image.ext}')
                with open(image_path, 'wb') as f:
                    f.write(image.blob)
                image_paths.append((slide_number, shape_number, image_path))
    
    return image_paths

def insert_alt_text(pptx_path, descriptions):
    logging.debug(f"Inserting alt text into {pptx_path}")
    prs = Presentation(pptx_path)
    
    for slide_number, shape_number, description in descriptions:
        logging.debug(f"Processing slide {slide_number}, shape {shape_number}")
        if 0 <= slide_number < len(prs.slides):
            slide = prs.slides[slide_number]
            if 0 <= shape_number < len(slide.shapes):
                shape = slide.shapes[shape_number]
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape._element.xpath('./p:nvPicPr/p:cNvPr')[0].set('descr', description)
                else:
                    logging.warning(f"Shape {shape_number} on slide {slide_number} is not a picture")
            else:
                logging.warning(f"Invalid shape number {shape_number} for slide {slide_number}")
        else:
            logging.warning(f"Invalid slide number {slide_number}")
    
    output_folder = current_app.config['OUTPUT_FOLDER']
    modified_pptx_path = os.path.join(output_folder, 'modified_' + os.path.basename(pptx_path))
    logging.info(f"Saving modified file to: {modified_pptx_path}")
    prs.save(modified_pptx_path)
    
    if os.path.exists(modified_pptx_path):
        logging.info(f"File saved successfully: {modified_pptx_path}")
    else:
        logging.error(f"Failed to save file: {modified_pptx_path}")
    
    return modified_pptx_path